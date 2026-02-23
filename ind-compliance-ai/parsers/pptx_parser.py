from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - optional runtime dependency
    Presentation = None  # type: ignore[assignment]

from parsers.common.atomic_fact_extractor import extract_atomic_facts


def parse_pptx(path: Path) -> dict[str, Any]:
    """Parse PPT/PPTX slides and speaker notes for supporting materials."""
    slides: list[dict[str, Any]] = []
    notes: list[dict[str, Any]] = []
    extracted_text: list[str] = []

    if path.suffix.lower() == ".pptx":
        if Presentation is None:
            raise RuntimeError("python-pptx is required for .pptx parsing. Install with: pip install python-pptx")
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides):
            slide_text_runs: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    text = str(shape.text).strip()
                    slide_text_runs.append(text)
                    extracted_text.append(text)

            note_text = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                note_text = slide.notes_slide.notes_text_frame.text.strip()
            if note_text:
                notes.append({"slide_number": slide_index + 1, "text": note_text})
                extracted_text.append(note_text)

            slides.append(
                {
                    "slide_number": slide_index + 1,
                    "text": "\n".join(slide_text_runs),
                }
            )
    else:
        decoded_text = path.read_bytes().decode("utf-8", errors="ignore")
        normalized_lines = [line.strip() for line in decoded_text.splitlines() if line.strip()]
        for index, line in enumerate(normalized_lines):
            slides.append({"slide_number": index + 1, "text": line})
        extracted_text.extend(normalized_lines)

    merged_text = "\n".join(extracted_text)
    ast_slides = [
        {
            "block_type": "slide",
            "slide_number": item["slide_number"],
            "text": item["text"],
            "notes": next(
                (note.get("text", "") for note in notes if note.get("slide_number") == item["slide_number"]),
                "",
            ),
        }
        for item in slides
    ]

    return {
        "source_path": str(path),
        "source_type": "presentation",
        "slides": slides,
        "notes": notes,
        "document_ast": {
            "source_type": "pptx" if path.suffix.lower() == ".pptx" else "ppt",
            "slides": ast_slides,
            "slide_count": len(ast_slides),
        },
        "text": merged_text,
        "atomic_facts": extract_atomic_facts(merged_text),
        "metadata": {
            "slide_count": len(slides),
            "note_count": len(notes),
            "parser_hint": "pptx-generic-v1",
        },
    }
